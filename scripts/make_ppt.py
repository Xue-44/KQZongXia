"""
市场推广PPT生成工具脚本
用法:
    python make_ppt.py --type market_analysis --data data.json --output report.pptx
    python make_ppt.py --type campaign_plan --output campaign.pptx
"""

import argparse
import json
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


# ============================================================
# 集团配色
# ============================================================
COLORS = {
    "primary": RGBColor(31, 56, 100),
    "secondary": RGBColor(68, 114, 196),
    "accent": RGBColor(237, 125, 49),
    "light": RGBColor(221, 235, 247),
    "dark": RGBColor(23, 42, 75),
    "text": RGBColor(51, 51, 51),
    "gray": RGBColor(128, 128, 128),
    "white": RGBColor(255, 255, 255),
    "green": RGBColor(0, 128, 0),
    "red": RGBColor(192, 0, 0),
}


def _add_top_bar(slide):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.1)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["primary"]
    bar.line.fill.background()


def _add_slide_title(slide, title):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.6))
    p = box.text_frame.add_paragraph()
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]


# ============================================================
# 页面构建
# ============================================================
def create_cover(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide)

    # 左侧装饰条
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.06), Inches(7.5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["accent"]
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(7.5), Inches(1.5))
    tf = tb.text_frame
    p = tf.add_paragraph()
    p.text = data.get("title", "市场推广分析报告")
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]

    sb = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(7.5), Inches(0.8))
    tf = sb.text_frame
    p = tf.add_paragraph()
    p.text = data.get("subtitle", "")
    p.font.size = Pt(22)
    p.font.color.rgb = COLORS["secondary"]

    ib = slide.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(7.5), Inches(1.2))
    tf = ib.text_frame
    for line in [
        f"报告日期: {data.get('date', datetime.now().strftime('%Y年%m月%d日'))}",
        f"汇报人: {data.get('presenter', '')}",
        f"部门: {data.get('dept', '')}",
        "内部资料，请勿外传",
    ]:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS["gray"]
        p.space_after = Pt(4)


def create_toc(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.7))
    p = tb.text_frame.add_paragraph()
    p.text = "目  录"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]

    sections = data.get("sections", [])
    for idx, sec in enumerate(sections, 1):
        y = Inches(1.5) + Inches(idx * 0.9)
        num = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(0.6), Inches(0.5)
        )
        num.fill.solid()
        num.fill.fore_color.rgb = COLORS["primary"]
        num.line.fill.background()
        p = num.text_frame.add_paragraph()
        p.text = f"0{idx}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS["white"]
        p.alignment = PP_ALIGN.CENTER

        cb = slide.shapes.add_textbox(Inches(1.6), y, Inches(6.5), Inches(0.5))
        p = cb.text_frame.add_paragraph()
        p.text = sec.get("title", "")
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS["text"]

        if sec.get("duration"):
            db = slide.shapes.add_textbox(Inches(8.2), y, Inches(1), Inches(0.5))
            p = db.text_frame.add_paragraph()
            p.text = sec["duration"]
            p.font.size = Pt(12)
            p.font.color.rgb = COLORS["gray"]

        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), y + Inches(0.7), Inches(8.4), Inches(0.01)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = COLORS["light"]
        line.line.fill.background()


def create_market_metrics(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide)
    _add_slide_title(slide, data.get("section_title", "市场环境概览"))

    metrics = data.get("metrics", [])
    card_w = Inches(2.1)
    gap = Inches(0.2)
    y = Inches(1.3)

    for idx, m in enumerate(metrics[:4]):
        x = Inches(0.5) + idx * (card_w + gap)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, Inches(1.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS["light"]
        card.line.fill.background()

        lb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1), card_w - Inches(0.3), Inches(0.3))
        p = lb.text_frame.add_paragraph()
        p.text = m.get("label", "")
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS["gray"]

        vb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.4), card_w - Inches(0.3), Inches(0.5))
        p = vb.text_frame.add_paragraph()
        p.text = str(m.get("value", ""))
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = COLORS["primary"]

        cb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(1.0), card_w - Inches(0.3), Inches(0.3))
        p = cb.text_frame.add_paragraph()
        change = m.get("change", "")
        p.text = change
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS["green"] if change.startswith("+") else (
            COLORS["red"] if change.startswith("-") else COLORS["gray"]
        )

    # 洞察
    insights = data.get("insights", [])
    if insights:
        ib = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(3.5))
        tf = ib.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = "核心洞察"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS["primary"]
        p.space_after = Pt(10)
        for ins in insights:
            p = tf.add_paragraph()
            p.text = f"  {chr(9654)} {ins}"
            p.font.size = Pt(13)
            p.font.color.rgb = COLORS["text"]
            p.space_after = Pt(8)


def create_table_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide)
    _add_slide_title(slide, data.get("section_title", "数据明细"))

    headers = data.get("headers", [])
    rows = data.get("rows", [])
    if not headers or not rows:
        return

    nr = len(rows) + 1
    nc = len(headers)
    tbl = slide.shapes.add_table(
        nr, nc, Inches(0.5), Inches(1.3), Inches(9), Inches(0.38) * min(nr, 14)
    ).table

    for ci, h in enumerate(headers):
        c = tbl.cell(0, ci)
        c.text = h
        c.fill.solid()
        c.fill.fore_color.rgb = COLORS["primary"]
        for p in c.text_frame.paragraphs:
            p.font.color.rgb = COLORS["white"]
            p.font.bold = True
            p.font.size = Pt(10)
            p.alignment = PP_ALIGN.CENTER

    for ri, row in enumerate(rows, 1):
        for ci, val in enumerate(row):
            c = tbl.cell(ri, ci)
            c.text = str(val)
            if ri % 2 == 1:
                c.fill.solid()
                c.fill.fore_color.rgb = COLORS["light"]
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.alignment = PP_ALIGN.CENTER


def create_budget_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide)
    _add_slide_title(slide, data.get("section_title", "费用执行报告"))

    # 左侧表格
    budget_rows = data.get("budget_data", [])
    if budget_rows:
        headers = ["项目", "预算(万)", "已执行(万)", "执行率", "剩余(万)"]
        nr = len(budget_rows) + 1
        tbl = slide.shapes.add_table(
            nr, 5, Inches(0.5), Inches(1.3), Inches(4.5), Inches(0.35) * min(nr, 14)
        ).table
        for ci, h in enumerate(headers):
            c = tbl.cell(0, ci)
            c.text = h
            c.fill.solid()
            c.fill.fore_color.rgb = COLORS["primary"]
            for p in c.text_frame.paragraphs:
                p.font.color.rgb = COLORS["white"]
                p.font.bold = True
                p.font.size = Pt(9)
                p.alignment = PP_ALIGN.CENTER
        for ri, row in enumerate(budget_rows, 1):
            for ci, val in enumerate(row):
                c = tbl.cell(ri, ci)
                c.text = str(val)
                for p in c.text_frame.paragraphs:
                    p.font.size = Pt(8)
                    p.alignment = PP_ALIGN.CENTER

    # 右侧费效分析
    fee_data = data.get("fee_ratio_data", [])
    if fee_data:
        rb = slide.shapes.add_textbox(Inches(5.5), Inches(1.3), Inches(4), Inches(5))
        tf = rb.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = "费效比分析"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS["primary"]
        p.space_after = Pt(12)
        for item in fee_data:
            p = tf.add_paragraph()
            p.text = f"  {item.get('label', '')}: {item.get('value', '')}"
            p.font.size = Pt(12)
            p.font.color.rgb = COLORS["text"]
            p.space_after = Pt(6)


def create_summary(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_top_bar(slide)
    _add_slide_title(slide, "总结与建议")

    box1 = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.2), Inches(4.5))
    tf = box1.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "核心结论"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]
    for pt in data.get("key_points", []):
        p = tf.add_paragraph()
        p.text = f"  {chr(9679)} {pt}"
        p.font.size = Pt(14)
        p.space_after = Pt(8)

    box2 = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(4.5))
    tf = box2.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = "行动建议"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]
    for idx, step in enumerate(data.get("next_steps", []), 1):
        p = tf.add_paragraph()
        p.text = f"  {idx}. {step}"
        p.font.size = Pt(14)
        p.space_after = Pt(8)

    # 风险提示
    risks = data.get("risks", [])
    if risks:
        box3 = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.5))
        tf = box3.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = f"  {chr(9888)} 风险提示: {', '.join(risks)}"
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS["red"]


# ============================================================
# 主流程
# ============================================================
TYPE_STRUCTURE = {
    "market_analysis": [create_cover, create_toc, create_market_metrics, create_table_slide, create_summary],
    "campaign_plan": [create_cover, create_toc, create_market_metrics, create_budget_slide, create_summary],
    "budget_report": [create_cover, create_budget_slide, create_table_slide, create_summary],
    "annual_plan": [create_cover, create_toc, create_market_metrics, create_table_slide, create_budget_slide, create_summary],
    "brand_strategy": [create_cover, create_toc, create_table_slide, create_market_metrics, create_summary],
}


def generate(ppt_type, data, output_path):
    if ppt_type not in TYPE_STRUCTURE:
        raise ValueError(f"不支持的类型: {ppt_type}")

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    sections = data.get("sections", [])
    handlers = TYPE_STRUCTURE[ppt_type]

    for i, handler in enumerate(handlers):
        section_data = sections[i] if i < len(sections) else {}
        merged = {**data, **section_data}
        handler(prs, merged)

    prs.save(output_path)
    return output_path


def main():
    parser = argparse.ArgumentParser(description="市场推广PPT生成工具")
    parser.add_argument("--type", dest="ppt_type", default="market_analysis",
                        choices=list(TYPE_STRUCTURE.keys()),
                        help="PPT类型")
    parser.add_argument("--data", help="JSON数据文件路径")
    parser.add_argument("--output", required=True, help="输出PPT文件路径")
    args = parser.parse_args()

    if args.data:
        with open(args.data, "r", encoding="utf-8") as f:
            data = json.load(f)
    else:
        data = {
            "title": "市场推广分析报告",
            "date": datetime.now().strftime("%Y年%m月%d日"),
            "sections": [],
        }

    output = generate(args.ppt_type, data, args.output)
    print(f"PPT 已生成: {output}")


if __name__ == "__main__":
    main()
